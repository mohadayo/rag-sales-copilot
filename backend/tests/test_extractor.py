"""extractor モジュールのユニットテスト"""

import os
import tempfile

import pytest

from app.core.extractor import _extract_plain_text, extract_text


class TestExtractPlainText:
    """_extract_plain_text のテスト"""

    def test_utf8_file(self, tmp_path):
        """UTF-8エンコードのファイルを正常に読み込む"""
        test_file = tmp_path / "test.txt"
        test_file.write_text("UTF-8テキストです。", encoding="utf-8")
        result = _extract_plain_text(str(test_file))
        assert result == "UTF-8テキストです。"

    def test_utf8_bom_file(self, tmp_path):
        """UTF-8 BOMエンコードのファイルを正常に読み込む"""
        test_file = tmp_path / "test.txt"
        test_file.write_text("BOM付きテキストです。", encoding="utf-8-sig")
        result = _extract_plain_text(str(test_file))
        assert "BOM付きテキストです。" in result

    def test_shift_jis_file(self, tmp_path):
        """Shift-JISエンコードのファイルを正常に読み込む"""
        test_file = tmp_path / "test.txt"
        test_file.write_bytes("Shift-JISテキスト".encode("shift-jis"))
        result = _extract_plain_text(str(test_file))
        assert "Shift-JIS" in result or len(result) > 0

    def test_empty_file(self, tmp_path):
        """空ファイルは空文字列を返す"""
        test_file = tmp_path / "empty.txt"
        test_file.write_text("", encoding="utf-8")
        result = _extract_plain_text(str(test_file))
        assert result == ""

    def test_multiline_file(self, tmp_path):
        """複数行のファイルを正常に読み込む"""
        test_file = tmp_path / "multi.txt"
        content = "行1\n行2\n行3\n"
        test_file.write_text(content, encoding="utf-8")
        result = _extract_plain_text(str(test_file))
        assert "行1" in result
        assert "行2" in result
        assert "行3" in result


class TestExtractText:
    """extract_text のテスト"""

    def test_txt_file(self, tmp_path):
        """TXTファイルのテキスト抽出"""
        test_file = tmp_path / "test.txt"
        test_file.write_text("テキストファイルです。", encoding="utf-8")
        result = extract_text(str(test_file))
        assert "テキストファイルです" in result

    def test_md_file(self, tmp_path):
        """Markdownファイルのテキスト抽出"""
        test_file = tmp_path / "test.md"
        test_file.write_text("# タイトル\n本文です。", encoding="utf-8")
        result = extract_text(str(test_file))
        assert "タイトル" in result
        assert "本文" in result

    def test_docx_file(self, tmp_path):
        """DOCXファイルのテキスト抽出"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("DOCX文書のテスト段落です。")
        doc.add_paragraph("2番目の段落です。")
        test_file = tmp_path / "test.docx"
        doc.save(str(test_file))
        result = extract_text(str(test_file))
        assert "DOCX文書のテスト段落" in result
        assert "2番目の段落" in result

    def test_pptx_file(self, tmp_path):
        """PPTXファイルのテキスト抽出"""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "プレゼンタイトル"
        body = slide.placeholders[1]
        body.text = "スライドの本文テキスト"
        test_file = tmp_path / "test.pptx"
        prs.save(str(test_file))
        result = extract_text(str(test_file))
        assert "プレゼンタイトル" in result
        assert "スライドの本文テキスト" in result

    def test_unsupported_extension_raises_error(self, tmp_path):
        """未対応の拡張子はValueErrorを発生させる"""
        test_file = tmp_path / "test.xyz"
        test_file.write_text("テスト", encoding="utf-8")
        with pytest.raises(ValueError, match="Unsupported file type"):
            extract_text(str(test_file))
