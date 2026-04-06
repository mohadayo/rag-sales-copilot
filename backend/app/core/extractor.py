"""テキスト抽出モジュール: PDF, DOCX, PPTX, TXT, Markdown に対応"""

import logging
import os

logger = logging.getLogger(__name__)


def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".doc": _extract_docx,
        ".pptx": _extract_pptx,
        ".txt": _extract_plain_text,
        ".md": _extract_plain_text,
    }
    extractor = extractors.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext}")
    logger.info("テキスト抽出開始: file_path=%s, ext=%s", file_path, ext)
    text = extractor(file_path)
    logger.info("テキスト抽出完了: file_path=%s, text_length=%d", file_path, len(text))
    return text


def _extract_pdf(file_path: str) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n".join(texts)


def _extract_docx(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


def _extract_plain_text(file_path: str) -> str:
    """テキスト・Markdownファイルを読み込む（エンコーディング自動検出）"""
    encodings = ["utf-8", "utf-8-sig", "shift-jis", "euc-jp", "latin-1"]
    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                content = f.read()
            logger.debug("テキストファイル読み込み成功: encoding=%s, file=%s", encoding, file_path)
            return content
        except UnicodeDecodeError:
            logger.debug("エンコーディング %s での読み込み失敗。次を試行: %s", encoding, file_path)
            continue
    # フォールバック: バイナリ読み込みで無効文字を置換
    logger.warning("全エンコーディング失敗。バイナリモードで読み込み: %s", file_path)
    with open(file_path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")
